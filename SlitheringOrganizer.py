import os
import json
import pickle
from pptx import Presentation
#C:\Users\Sushant\Dropbox\College\Summer 17\Python\SlitheringMover

#figure out how to organize college lecture pdf

#exclude when slithering
exclude=["Slither_organized","Folders","Downloadsdesktop.ini","desktop.ini"]

#ppt and pdf handler
def get_slidepage(loc):
	text_dict ={}
	slide_no=1
	slide_text=""
	prs = Presentation(loc)
	for slide in prs.slides:
		for shape in slide.shapes:
			if not shape.has_text_frame:
				continue
			for paragraph in shape.text_frame.paragraphs:
				for run in paragraph.runs:
					slide_text+=" "+run.text		
		slide_no+=1
		if slide_no>2 :
			break
	return slide_text.lower()

def get_course(course_names,pptxloc):
	info=get_slidepage(pptxloc)
	words=info.split()
	for course in course_names:
		if course in words:
			return course 
	return -1

def getallfiles(loc):
	files= os.listdir(loc);
	file_dic={}
	for file in files:
		if file in exclude: 
			continue;
		file_dic[file]=loc+"\\"+file
	return file_dic 

def file_extension_split(ext,file_dic):
	#file_dic= getallfiles(loc)
	ext_dic={}
	for file,file_loc in file_dic.iteritems():
		if file[-4:].lower() in ext :
			ext_dic[file]=file_loc
		if file[-5:].lower() in ext :
			ext_dic[file]=file_loc
	return ext_dic

def no_extension(file_dic):
	noext={}
	for file,file_loc in file_dic.iteritems():
		pos=file.find('.')
		if pos == -1:
			noext[file]=file_loc
	return noext

def moveto(og,newdir):
	for file,ogfile_loc in og.iteritems():
		try:
			#print ogfile_loc
			#print "\n"
			#print newdir+file
			os.rename(ogfile_loc,newdir+file)
		except WindowsError:
			pos=file.find('.')
			newfile=file[:pos]+"_dupl"+file[pos:]
			os.rename(ogfile_loc,newdir+newfile)


def load_obj(name ):
    with open(name + '.pkl', 'rb') as f:
        return pickle.load(f)

def organize():
	down_dir=""
	others_dir=""
	fromfiles = {}
	college_dir=""
	allfiles={}
	#open the file
	#fromfiles= load_obj("direc_and_ext")
	with open("direc_and_ext.txt") as file:
		fromfiles = json.load(file)
	
	#get names of courses
	course_names=fromfiles["Courses"]

	#get downloads directory
	for direc,key in fromfiles.iteritems():
		if key == "downloads":
			down_dir=direc
		if key == "others":
			others_dir=direc 
		if ".pdf" in key:  #assumes pdf to be in key
			college_dir=direc
	down_dic=getallfiles(down_dir)

	#forming nested dictionary based on extension
	for direc,key in fromfiles.iteritems():
		if key == "downloads" or key == "others" or direc=="Courses":
			continue 
		if key == "folders" :
			allfiles[direc]=no_extension(down_dic)
		else:
			allfiles[direc]=file_extension_split(key,down_dic)

	#handle pptx and pdf files
	college_files=allfiles[college_dir]
	for file,loc in college_files.iteritems():
		temp_dir=college_dir
		if file.find('.pptx')!=-1:
			course=get_course(course_names,loc)
			if course == -1:
				continue
			temp_dir+=course+"\\"
			if not os.path.exists(temp_dir):
				os.makedirs(temp_dir)
			moveto({file:loc},temp_dir)

	#update for remaining files
	down_dic=getallfiles(down_dir)
	allfiles[college_dir]=file_extension_split(fromfiles[college_dir],down_dic)
	
	#moving for most extension
	for dir_file,ext_dic in allfiles.iteritems():
		moveto(ext_dic,dir_file) 
	
	#for moving others
	down_dic=getallfiles(down_dir)
	moveto(down_dic,others_dir)
